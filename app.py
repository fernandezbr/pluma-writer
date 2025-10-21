import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

# --- NEW imports ---
from datetime import datetime
from io import BytesIO

# PDF
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os

# --- UI helper: centered "OR" header with lines ---
def or_header(text: str):
    st.markdown(
        """
        <style>
        .or-header{display:flex;align-items:center;gap:.75rem;margin:.5rem 0 1rem;}
        .or-header:before,.or-header:after{content:"";flex:1;border-top:1px solid rgba(128,128,128,.35);}
        .or-header .or-text{font-weight:600;opacity:.85;}
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(f"<div class='or-header'><span class='or-text'>{text}</span></div>", unsafe_allow_html=True)

# App title
pages.show_home()
pages.show_sidebar()

# Home page
st.header("📝Style Writer")

# # Content input
# st.session_state.content = st.text_area(
#     ":blue[**Content Input for BSP Writing Style:**]", st.session_state.content, 200
# )

# uploaded_files = st.file_uploader(
#     ":blue[**Upload Content Files:**]", 
#     type=["pdf", "docx", "pptx"], 
#     accept_multiple_files=True,
#     help="Upload PDF, Word, or PowerPoint files"
# )
or_header("Input the Contents or Upload the File for BSP Style Writing")

# --- two-column layout (Col 1 / Col 2) ---
col1, col2 = st.columns([3, 2], gap="small")

with col1:
    
    st.session_state.content = st.text_area(
        ":blue[**Input Content:**]",
        st.session_state.content,
        height=160,
        key="content_input",
    )

    max_output_length = st.slider(
        ":blue[**Output Maximum Character Length:**]", 
        min_value=20, 
        max_value=75000, 
        value=None, 
        step=None, 
        format=None, 
        key=None, 
        help=None, 
        on_change=None, 
        disabled=False, 
        label_visibility="visible", 
        width="stretch"
    )

with col2:
    
    uploaded_files = st.file_uploader(
        ":blue[**Upload Files:**]",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files",
        key="content_upload",
    )

# Extract text from uploaded files
extracted_text = ""
if uploaded_files:
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split('.')[-1].lower()
        
        if file_type == 'pdf':
            pdf_reader = PdfReader(BytesIO(uploaded_file.read()))
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() + "\n"
                
        elif file_type == 'docx':
            doc = Document(BytesIO(uploaded_file.read()))
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    extracted_text += paragraph.text + "\n"
                    
        elif file_type == 'pptx':
            prs = Presentation(BytesIO(uploaded_file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            extracted_text += shape.text + "\n"


# ----------------------------
# PICK EXACTLY ONE CONTENT SOURCE + PREVIEW
# ----------------------------
has_input = bool(st.session_state.content.strip())
has_uploads = bool(extracted_text.strip())

# Let the user choose if both are present; otherwise auto-pick
if has_input and has_uploads:
    source = st.radio(
        ":blue[**Choose content source**]",
        ["Uploaded files", "Manual input"],
        horizontal=True,
        index=0,
        help="Use either the text you typed OR the text extracted from the uploaded files."
    )
elif has_uploads:
    source = "Uploaded files"
elif has_input:
    source = "Manual input"
else:
    source = None

# Decide content_all and show a preview when using uploads
if source == "Uploaded files":
    content_all = extracted_text  # keep full unicode; your PDF builder handles fonts
    with st.expander("📄 Preview: Extracted text from uploaded files", expanded=True):
        st.text_area(
            "Extracted Text",
            content_all,
            height=240,
            key="uploaded_text_preview",
        )
elif source == "Manual input":
    content_all = st.session_state.content
else:
    content_all = ""
    st.markdown(
    """
    <div class="bsp-alert-red" role="alert">
      <strong>Heads up:</strong> Provide content — either type in the left box or upload a file on the right.
    </div>
    <style>
      .bsp-alert-red{
        padding:12px 14px;
        margin: 4px 0 10px;
        border-radius:10px;
        border:1px solid rgba(220,53,69,.35);
        background: rgba(220,53,69,.08); /* light red */
        font-size: 0.95rem;
      }
      .bsp-alert-red strong{
        color:#b02a37; /* dark red for emphasis */
      }
    </style>
    """,
    unsafe_allow_html=True,
)

# Combine text area and extracted content
#content_all = st.session_state.content + "\n" + extracted_text.encode("ascii", errors="ignore").decode("ascii")

# Extracting the styles and creating combined display options
styles_data = utils.get_styles()
style_options = [item['name'] for item in styles_data]
selected_style = st.selectbox(":blue[**Select a Style:**]", options=style_options, index=None)

# Assigning the selected style to the session state
if selected_style:
    # Find the matching style data
    filtered = next(
        (item for item in styles_data if str(item["name"]) == selected_style), None
    )
    
    if filtered:
        st.session_state.style = filtered["style"]
        st.session_state.example = filtered["example"]
        st.session_state.styleId = selected_style
        
# st.session_state.style = st.text_area(":blue[**Style:**]", st.session_state.style)

# Show the example style
guidelines = st.session_state.locals.get("relevant_guidelines", {})
guidelines_summary = st.session_state.locals.get("guideline_summaries", {}) 
selected_guidelines = []

st.write(":blue[**Select Editorial Style Guides:**]")

# Tooltip for guideline summary in the UI
def render_guideline_checkbox(section_name: str, content: str, col_key_prefix: str):
    default_checked = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
    tooltip = guidelines_summary.get(section_name, None)  # one-sentence summary for hover
    if st.checkbox(
        section_name,
        value=default_checked,
        key=f"{col_key_prefix}_{section_name}",
        help=tooltip  # <-- hover tooltip appears on the ⓘ icon and on hover
    ):
        selected_guidelines.append(content)

# Create a checkbox for each guideline section
if guidelines:
    with st.container(border=True):
        # Create two columns
        col1, col2 = st.columns(2)

        # Split guidelines into two halves
        guideline_items = list(guidelines.items())
        mid_point = len(guideline_items) // 2

        # First column
        with col1:
            for section_name, content in guideline_items[:mid_point]:
                render_guideline_checkbox(section_name, content, "col1")
                # default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
                # if st.checkbox(section_name, value=default, key=f"col1_{section_name}"):
                #     selected_guidelines.append(content)

        # Second column
        with col2:
            for section_name, content in guideline_items[mid_point:]:
                render_guideline_checkbox(section_name, content, "col2")
                # default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
                # if st.checkbox(section_name, value=default, key=f"col2_{section_name}"):
                #     selected_guidelines.append(content)
else:
    st.warning("No guidelines available in the local data.")

# Join all selected guidelines with newlines and store in session state
st.session_state.guidelines = "\n".join(selected_guidelines)

# Show the combined guidelines in a text area
# st.text_area(":blue[**Relevant Guidelines:**]", st.session_state.guidelines, height=200)


# --- NEW helpers: make DOCX/PDF from text ---

def make_docx_bytes(text: str, title: str | None = None) -> bytes:
    """Return a .docx file (bytes) with a title and body paragraphs."""
    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    # Split into paragraphs on blank lines while preserving line breaks
    for block in text.replace("\r\n", "\n").split("\n\n"):
        p = doc.add_paragraph()
        for line in block.split("\n"):
            if line.strip():
                p.add_run(line)
            p.add_run("\n")
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _register_pdf_font_if_available():
    """Optionally register DejaVuSans for better Unicode PDF rendering."""
    try:
        font_path = os.path.join("assets", "DejaVuSans.ttf")
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont("DejaVuSans", font_path))
            return "DejaVuSans"
    except Exception:
        pass
    # Fallback to built-in Helvetica (ASCII/Latin-1 safe)
    return "Helvetica"


def make_pdf_bytes(text: str, title: str | None = None) -> bytes:
    """Return a PDF (bytes) using ReportLab."""
    font_name = _register_pdf_font_if_available()

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        title=title or "Rewrite",
        author="Style Writer",
    )

    styles = getSampleStyleSheet()
    base = styles["BodyText"]
    base.fontName = font_name
    base.fontSize = 11
    base.leading = 14

    title_style = ParagraphStyle(
        "Title",
        parent=styles["Heading1"],
        fontName=font_name,
        spaceAfter=12,
    )

    story = []
    if title:
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 8))

    # Turn double newlines into paragraph breaks; single newlines stay inside a paragraph
    for block in text.replace("\r\n", "\n").split("\n\n"):
        # Escape simple HTML-sensitive chars for Paragraph
        block = block.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        block = block.replace("\n", "<br/>")
        story.append(Paragraph(block, base))
        story.append(Spacer(1, 6))

    doc.build(story)
    return buf.getvalue()


if st.button(
    ":blue[**Rewrite Content**]",
    key="extract",
   # disabled=content_all == ""
    disabled=(
    content_all.strip() == ""
    or st.session_state.style == ""
    or st.session_state.example == ""
)
    or st.session_state.style == ""
    or st.session_state.example == "",
):
    with st.container(border=True):
        with st.spinner("Processing..."):
            # --- NEW: show the result and download buttons ---
            st.markdown("### ✨ Rewritten Output")
            output = prompts.rewrite_content(content_all, max_output_length, False)
            utils.save_output(output, content_all)

            # --- NEW: cache for later & build filenames ---
            st.session_state["last_output"] = output
            ts = datetime.now().strftime("%Y%m%d-%H%M%S")
            style_id = (st.session_state.get("styleId") or "Style").replace(" ", "_")
            base_name = f"rewrite_{style_id}_{ts}"

            # --- NEW: build bytes for DOCX and PDF ---
            title_text = f"Rewrite • {st.session_state.get('styleId') or 'Selected Style'}"
            docx_bytes = make_docx_bytes(output, title=title_text)
            pdf_bytes = make_pdf_bytes(output, title=title_text)

           # st.text_area("Result", output, height=300)

            c1, c2 = st.columns(2)
            with c1:
                st.download_button(
                    "⬇️ Download as DOCX",
                    data=docx_bytes,
                    file_name=f"{base_name}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                )
            with c2:
                st.download_button(
                    "⬇️ Download as PDF",
                    data=pdf_bytes,
                    file_name=f"{base_name}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                )